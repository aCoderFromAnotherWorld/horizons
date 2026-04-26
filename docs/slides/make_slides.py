"""
Horizons: ASD Pre-Screening Serious Game — Presentation Generator
CSE 3200: System Development Project | KUET, April 2026
Authors: Md. Shifat Hasan (2107067) & Abu Hasanat Soykot (2107100)
Supervisor: Dr. Md. Aminul Haque Akhand
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ─── Palette ────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

OFF_WHITE    = rgb(0xFA, 0xFA, 0xF8)   # background
DEEP_INDIGO  = rgb(0x1E, 0x1B, 0x4B)  # footer, heavy headings
INDIGO_700   = rgb(0x43, 0x38, 0xCA)  # left strip, strong accent
INDIGO_600   = rgb(0x49, 0x3C, 0xD1)
INDIGO_100   = rgb(0xE0, 0xE7, 0xFF)  # header bar
INDIGO_200   = rgb(0xC7, 0xD2, 0xFE)
INDIGO_50    = rgb(0xEE, 0xF2, 0xFF)
SLATE_800    = rgb(0x1E, 0x29, 0x3B)  # body text
SLATE_600    = rgb(0x47, 0x55, 0x69)  # secondary text
AMBER_600    = rgb(0xD9, 0x77, 0x06)  # red-flag accent
AMBER_100    = rgb(0xFE, 0xF3, 0xC7)
GREEN_700    = rgb(0x04, 0x78, 0x57)  # positive accent
GREEN_100    = rgb(0xD1, 0xFA, 0xE5)
RED_600      = rgb(0xDC, 0x26, 0x26)
WHITE        = rgb(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = rgb(0xF1, 0xF5, 0xF9)
LAVENDER     = rgb(0xA5, 0xB4, 0xFC)  # footer text

# ─── Slide dimensions (10" × 5.625" widescreen) ─────────────────────────────
W = Inches(10)
H = Inches(5.625)

LEFT_STRIP_W  = Inches(0.19)
HEADER_H      = Inches(0.72)
FOOTER_H      = Inches(0.52)

BODY_X = LEFT_STRIP_W + Inches(0.25)
BODY_W = W - BODY_X - Inches(0.25)
BODY_Y_START  = HEADER_H + Inches(0.18)
BODY_Y_END    = H - FOOTER_H - Inches(0.12)
BODY_H        = BODY_Y_END - BODY_Y_START

FOOTER_TEXT = (
    "Horizons: A Serious Game for ASD Pre-Screening   |   "
    "Md. Shifat Hasan & Abu Hasanat Soykot   |   KUET, 2026"
)
TOTAL_SLIDES = 12


# ─── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txb(slide, text, x, y, w, h,
        size=16, bold=False, italic=False,
        color=SLATE_800, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_para(tf, text, size=15, bold=False, italic=False,
             color=SLATE_800, align=PP_ALIGN.LEFT, space_before_pt=0):
    para = tf.add_paragraph()
    para.alignment = align
    if space_before_pt:
        para.space_before = Pt(space_before_pt)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return para


def make_base_slide(prs, title, slide_num, total=TOTAL_SLIDES):
    """Creates a slide with the standard chrome (left strip, header, footer)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Off-white background
    add_rect(slide, 0, 0, W, H, OFF_WHITE)

    # Left accent strip (full height)
    add_rect(slide, 0, 0, LEFT_STRIP_W, H, INDIGO_700)

    # Header bar
    add_rect(slide, 0, 0, W, HEADER_H, INDIGO_100)
    add_rect(slide, 0, 0, LEFT_STRIP_W, HEADER_H, INDIGO_700)

    # Header title text
    txb(slide, title,
        BODY_X, Inches(0.13),
        BODY_W, HEADER_H - Inches(0.13),
        size=22, bold=True, color=DEEP_INDIGO, align=PP_ALIGN.LEFT)

    # Footer bar
    footer_y = H - FOOTER_H
    add_rect(slide, 0, footer_y, W, FOOTER_H, DEEP_INDIGO)
    add_rect(slide, 0, footer_y, LEFT_STRIP_W, FOOTER_H, INDIGO_700)

    # Footer running text
    txb(slide, FOOTER_TEXT,
        BODY_X, footer_y + Inches(0.09),
        W - LEFT_STRIP_W - Inches(1.4), FOOTER_H - Inches(0.09),
        size=8.5, color=LAVENDER, align=PP_ALIGN.LEFT)

    # Slide number
    txb(slide, f"{slide_num} / {total}",
        W - Inches(1.3), footer_y + Inches(0.09),
        Inches(1.15), FOOTER_H - Inches(0.09),
        size=8.5, bold=True, color=LAVENDER, align=PP_ALIGN.RIGHT)

    return slide


def bullet_block(slide, items, x, y, w,
                 marker="▸", marker_color=INDIGO_700,
                 size=15.5, bold=False, color=SLATE_800,
                 line_gap=Inches(0.37)):
    """
    items: list of strings (or (text, sub_items) tuples)
    Returns the y-position after the last bullet.
    """
    cur_y = y
    for item in items:
        # marker
        txb(slide, marker,
            x, cur_y, Inches(0.28), line_gap,
            size=size, bold=True, color=marker_color)
        # text
        txb(slide, item,
            x + Inches(0.28), cur_y, w - Inches(0.28), line_gap,
            size=size, bold=bold, color=color, wrap=True)
        cur_y += line_gap
    return cur_y


def info_box(slide, label, value, x, y, w, h,
             bg=INDIGO_50, border=INDIGO_200,
             label_size=10, val_size=22, val_bold=True):
    """Small stat/info card."""
    add_rect(slide, x, y, w, h, bg, line_color=border, line_width=Pt(0.75))
    txb(slide, label, x + Inches(0.1), y + Inches(0.06),
        w - Inches(0.2), Inches(0.3),
        size=label_size, color=SLATE_600, align=PP_ALIGN.CENTER)
    txb(slide, value, x + Inches(0.1), y + Inches(0.3),
        w - Inches(0.2), h - Inches(0.35),
        size=val_size, bold=val_bold, color=DEEP_INDIGO, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Full off-white background
    add_rect(slide, 0, 0, W, H, OFF_WHITE)

    # Left decorative strip (wide for title page)
    add_rect(slide, 0, 0, Inches(0.45), H, INDIGO_700)

    # Bottom accent bar
    add_rect(slide, 0, H - Inches(0.70), W, Inches(0.70), DEEP_INDIGO)
    add_rect(slide, 0, H - Inches(0.70), Inches(0.45), Inches(0.70), INDIGO_700)

    # Thin horizontal divider mid-slide
    add_rect(slide, Inches(0.65), Inches(2.60), W - Inches(0.9), Inches(0.025), INDIGO_200)

    # Course label
    txb(slide, "CSE 3200: System Development Project",
        Inches(0.65), Inches(0.35),
        W - Inches(0.9), Inches(0.45),
        size=12, italic=True, color=SLATE_600, align=PP_ALIGN.LEFT)

    # Main title
    txb(slide, "HORIZONS",
        Inches(0.65), Inches(0.80),
        W - Inches(0.9), Inches(1.10),
        size=48, bold=True, color=INDIGO_700, align=PP_ALIGN.LEFT)

    # Subtitle
    txb(slide, "A Serious Game for Autism Spectrum Disorder Pre-Screening",
        Inches(0.65), Inches(1.88),
        W - Inches(0.9), Inches(0.70),
        size=19, bold=False, color=SLATE_800, align=PP_ALIGN.LEFT)

    # Divider already drawn

    # Authors + supervisor block
    txb(slide, "Md. Shifat Hasan  (Roll: 2107067)   &   Abu Hasanat Soykot  (Roll: 2107100)",
        Inches(0.65), Inches(2.75),
        W - Inches(0.9), Inches(0.42),
        size=14, bold=True, color=DEEP_INDIGO, align=PP_ALIGN.LEFT)

    txb(slide, "Supervisor: Dr. Md. Aminul Haque Akhand, Professor, Dept. of CSE, KUET",
        Inches(0.65), Inches(3.18),
        W - Inches(0.9), Inches(0.38),
        size=12.5, italic=True, color=SLATE_600, align=PP_ALIGN.LEFT)

    # Institution
    txb(slide, "Department of Computer Science and Engineering\n"
               "Khulna University of Engineering & Technology, Khulna 9203, Bangladesh",
        Inches(0.65), Inches(3.60),
        W - Inches(0.9), Inches(0.65),
        size=12, color=SLATE_800, align=PP_ALIGN.LEFT)

    # Date in footer
    txb(slide, "April 2026",
        Inches(0.65), H - Inches(0.57),
        W - Inches(0.9), Inches(0.42),
        size=11, color=LAVENDER, align=PP_ALIGN.LEFT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — OUTLINE
# ═══════════════════════════════════════════════════════════════════════════
def slide_outline(prs):
    slide = make_base_slide(prs, "Outline", 1)

    items_left = [
        "1.  The Problem — ASD screening gap",
        "2.  State of the Art & Research Gaps",
        "3.  Horizons — What We Built",
        "4.  System Architecture",
        "5.  Game Design: 6 Chapters",
    ]
    items_right = [
        "6.  Scoring Engine & Red-Flag Detection",
        "7.  Implementation Results",
        "8.  Ethical Guardrails & Impact",
        "9.  Conclusion & Future Work",
    ]

    col_w = (BODY_W - Inches(0.2)) / 2
    col2_x = BODY_X + col_w + Inches(0.2)

    cur_y = BODY_Y_START + Inches(0.08)
    for item in items_left:
        txb(slide, item, BODY_X, cur_y, col_w, Inches(0.40),
            size=15.5, color=SLATE_800, align=PP_ALIGN.LEFT)
        cur_y += Inches(0.42)

    cur_y = BODY_Y_START + Inches(0.08)
    for item in items_right:
        txb(slide, item, col2_x, cur_y, col_w, Inches(0.40),
            size=15.5, color=SLATE_800, align=PP_ALIGN.LEFT)
        cur_y += Inches(0.42)

    # Vertical divider between columns
    add_rect(slide, BODY_X + col_w + Inches(0.08), BODY_Y_START,
             Inches(0.025), BODY_H, INDIGO_200)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
def slide_problem(prs):
    slide = make_base_slide(prs, "The Problem — ASD Screening Gap", 2)

    # Key stat cards across the top
    card_w = Inches(2.55)
    card_h = Inches(1.08)
    card_y = BODY_Y_START + Inches(0.05)
    gap = Inches(0.20)
    cards = [
        ("1 in 36 children\n(USA)", "Maenner et al., 2023"),
        ("1–2% of children\nglobally", "Baxter et al., 2015"),
        ("Avg. diagnosis\ndelay: 2–4 years", "Under-resourced settings"),
    ]
    cx = BODY_X
    for val, lbl in cards:
        add_rect(slide, cx, card_y, card_w, card_h, INDIGO_50, INDIGO_200, Pt(0.75))
        txb(slide, val, cx + Inches(0.1), card_y + Inches(0.1),
            card_w - Inches(0.2), card_h - Inches(0.4),
            size=16.5, bold=True, color=DEEP_INDIGO, align=PP_ALIGN.CENTER)
        txb(slide, lbl, cx + Inches(0.1), card_y + card_h - Inches(0.32),
            card_w - Inches(0.2), Inches(0.28),
            size=9, italic=True, color=SLATE_600, align=PP_ALIGN.CENTER)
        cx += card_w + gap

    # Body bullets
    bul_y = card_y + card_h + Inches(0.22)
    bullets = [
        "Gold standard: ADOS-2 — requires trained clinicians, controlled environment, and considerable cost",
        "Caregiver questionnaires (M-CHAT-R/F) — practical but vulnerable to observer recall bias",
        "Serious games emerged as a middle ground — but existing tools are narrow in age range,\n"
        "  proprietary, install-required, or single-domain only",
    ]

    cur_y = bul_y
    for b in bullets:
        txb(slide, "▸", BODY_X, cur_y, Inches(0.28), Inches(0.50),
            size=15, bold=True, color=INDIGO_700)
        txb(slide, b, BODY_X + Inches(0.28), cur_y,
            BODY_W - Inches(0.28), Inches(0.52),
            size=14.5, color=SLATE_800, wrap=True)
        cur_y += Inches(0.56)

    # Research question box
    rq_y = cur_y + Inches(0.08)
    add_rect(slide, BODY_X, rq_y, BODY_W, Inches(0.58), AMBER_100,
             AMBER_600, Pt(0.75))
    txb(slide, "How can a browser-delivered serious game provide engaging experiences for children aged 3–10 "
               "and collect clinically meaningful, multi-domain ASD screening signals — without specialist equipment?",
        BODY_X + Inches(0.15), rq_y + Inches(0.07),
        BODY_W - Inches(0.3), Inches(0.50),
        size=12.5, italic=True, color=rgb(0x78, 0x35, 0x00), align=PP_ALIGN.LEFT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — STATE OF THE ART & GAPS
# ═══════════════════════════════════════════════════════════════════════════
def slide_related_work(prs):
    slide = make_base_slide(prs, "State of the Art & Research Gaps", 3)

    # Comparison table
    col_labels = ["Feature", "EmoGalaxy", "DTT Game", "ADOS-2", "Horizons"]
    rows = [
        ["Browser-based",     "✗", "✗", "✗", "✓"],
        ["Multi-domain (all 4)", "✗", "✗", "✗", "✓"],
        ["Age range 3–10",    "6–14", "ASD only", "2–16", "3–10"],
        ["Automated scoring", "✓", "✓", "✗", "✓"],
        ["Red-flag detection","✗", "✗", "✓", "✓"],
        ["Pointer heatmap",   "✗", "✗", "✗", "✓"],
        ["Caregiver report",  "✗", "✗", "✗", "✓"],
        ["Open source",       "✗", "✗", "✗", "✓"],
    ]

    th_bg     = DEEP_INDIGO
    row_bg    = [OFF_WHITE, LIGHT_GRAY]
    hor_bg    = INDIGO_100
    hor_col   = INDIGO_700

    col_widths = [Inches(2.1), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.35)]
    row_h = Inches(0.37)
    tbl_x = BODY_X
    tbl_y = BODY_Y_START + Inches(0.05)

    # Header row
    cx = tbl_x
    for j, lbl in enumerate(col_labels):
        is_horizons = (j == 4)
        bg = hor_bg if is_horizons else th_bg
        tc = hor_col if is_horizons else WHITE
        add_rect(slide, cx, tbl_y, col_widths[j], row_h, bg)
        txb(slide, lbl, cx + Inches(0.05), tbl_y + Inches(0.06),
            col_widths[j] - Inches(0.1), row_h - Inches(0.06),
            size=12, bold=True, color=tc, align=PP_ALIGN.CENTER)
        cx += col_widths[j]

    # Data rows
    for i, row in enumerate(rows):
        cur_y = tbl_y + row_h * (i + 1)
        cx = tbl_x
        for j, cell in enumerate(row):
            is_horizons = (j == 4)
            bg = hor_bg if is_horizons else row_bg[i % 2]
            is_yes = cell == "✓"
            is_no  = cell == "✗"
            col = hor_col if is_horizons else (GREEN_700 if is_yes else (RED_600 if is_no else SLATE_800))
            bld = is_yes or is_no or is_horizons
            add_rect(slide, cx, cur_y, col_widths[j], row_h, bg)
            txb(slide, cell, cx + Inches(0.05), cur_y + Inches(0.06),
                col_widths[j] - Inches(0.1), row_h - Inches(0.06),
                size=12.5, bold=bld, color=col, align=PP_ALIGN.CENTER)
            cx += col_widths[j]

    # Key insight note below table
    note_y = tbl_y + row_h * (len(rows) + 1) + Inches(0.12)
    txb(slide, "Horizons is the only entry simultaneously offering browser delivery, multi-domain coverage, "
               "red-flag detection, and a shareable caregiver report — while remaining fully open source.",
        BODY_X, note_y, BODY_W, Inches(0.42),
        size=12, italic=True, color=SLATE_600, align=PP_ALIGN.LEFT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — HORIZONS OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
def slide_overview(prs):
    slide = make_base_slide(prs, "Horizons — What We Built", 4)

    # Two-column: left = what it is / right = novelty
    col_w = (BODY_W - Inches(0.25)) / 2
    col2_x = BODY_X + col_w + Inches(0.25)

    # Left column header
    add_rect(slide, BODY_X, BODY_Y_START, col_w, Inches(0.38), INDIGO_100)
    txb(slide, "The System",
        BODY_X + Inches(0.1), BODY_Y_START + Inches(0.04),
        col_w - Inches(0.2), Inches(0.30),
        size=14, bold=True, color=DEEP_INDIGO)

    left_items = [
        "Browser-based serious game — no installation, no dedicated hardware",
        "6 interconnected gameplay chapters, ages 3–10",
        "Covers all 4 ASD indicator domains (ADOS-2 / DSM-5 / M-CHAT-R/F)",
        "Records: response time · accuracy · pointer trajectory · break count",
        "HMAC-SHA256 signed caregiver report at session end",
        "Researcher dashboard: domain charts · heatmaps · CSV/JSON export",
    ]
    cur_y = BODY_Y_START + Inches(0.42)
    for item in left_items:
        txb(slide, "▸", BODY_X, cur_y, Inches(0.26), Inches(0.44),
            size=14, bold=True, color=INDIGO_700)
        txb(slide, item, BODY_X + Inches(0.26), cur_y,
            col_w - Inches(0.26), Inches(0.46),
            size=13, color=SLATE_800, wrap=True)
        cur_y += Inches(0.50)

    # Right column header
    add_rect(slide, col2_x, BODY_Y_START, col_w, Inches(0.38), INDIGO_100)
    txb(slide, "What Makes It Novel",
        col2_x + Inches(0.1), BODY_Y_START + Inches(0.04),
        col_w - Inches(0.2), Inches(0.30),
        size=14, bold=True, color=DEEP_INDIGO)

    right_items = [
        "Only tool combining all four ADOS-2 / DSM-5 domains in one browser session",
        "Penalty scoring grounded in 3 independent peer-reviewed frameworks",
        "5 automatic red-flag detectors with multiplicative severity weighting",
        "Real-time audio synthesis via Tone.js — zero audio files shipped",
        "312 automated tests — scoring, DB queries, API routes, full pipeline",
        "Fully open source · deployable on Vercel free tier",
    ]
    cur_y = BODY_Y_START + Inches(0.42)
    for item in right_items:
        txb(slide, "▸", col2_x, cur_y, Inches(0.26), Inches(0.44),
            size=14, bold=True, color=AMBER_600)
        txb(slide, item, col2_x + Inches(0.26), cur_y,
            col_w - Inches(0.26), Inches(0.46),
            size=13, color=SLATE_800, wrap=True)
        cur_y += Inches(0.50)

    # Vertical divider
    mid_x = BODY_X + col_w + Inches(0.10)
    add_rect(slide, mid_x, BODY_Y_START, Inches(0.025), BODY_H, INDIGO_200)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
def slide_architecture(prs):
    slide = make_base_slide(prs, "System Architecture", 5)

    # 4-tier block diagram
    tier_data = [
        ("Browser Tier",
         "React · Zustand · Framer Motion · Tone.js · @dnd-kit",
         INDIGO_50, INDIGO_700),
        ("Application Tier",
         "Next.js 15 App Router · Scoring Engine · Better Auth · HMAC Token",
         GREEN_100, GREEN_700),
        ("Data Tier",
         "PostgreSQL via Neon · postgres.js · Auto-migration (schema.sql)",
         AMBER_100, AMBER_600),
        ("External Services",
         "Vercel Hosting · Resend Email · Vercel Analytics",
         LIGHT_GRAY, SLATE_600),
    ]

    block_w = Inches(2.1)
    block_h = Inches(1.12)
    arrow_w = Inches(0.28)
    total_w = len(tier_data) * block_w + (len(tier_data) - 1) * arrow_w
    start_x = BODY_X + (BODY_W - total_w) / 2
    block_y = BODY_Y_START + Inches(0.12)

    cx = start_x
    for i, (title, subtitle, bg, col) in enumerate(tier_data):
        add_rect(slide, cx, block_y, block_w, block_h, bg, col, Pt(1.5))
        txb(slide, title,
            cx + Inches(0.1), block_y + Inches(0.1),
            block_w - Inches(0.2), Inches(0.38),
            size=13.5, bold=True, color=col, align=PP_ALIGN.CENTER)
        txb(slide, subtitle,
            cx + Inches(0.1), block_y + Inches(0.48),
            block_w - Inches(0.2), block_h - Inches(0.52),
            size=10, color=col, align=PP_ALIGN.CENTER, wrap=True)
        if i < len(tier_data) - 1:
            ax = cx + block_w
            ay = block_y + block_h / 2 - Inches(0.12)
            txb(slide, "→", ax, ay, arrow_w, Inches(0.28),
                size=16, bold=True, color=SLATE_600, align=PP_ALIGN.CENTER)
        cx += block_w + arrow_w

    # Key data flows label
    txb(slide, "HTTP/JSON · postgres.js / TLS · HTTPS API",
        start_x, block_y + block_h + Inches(0.08),
        total_w, Inches(0.30),
        size=10.5, italic=True, color=SLATE_600, align=PP_ALIGN.CENTER)

    # Separator
    sep_y = block_y + block_h + Inches(0.46)
    add_rect(slide, BODY_X, sep_y, BODY_W, Inches(0.022), INDIGO_200)

    # Data collected row
    txb(slide, "Data Collected Per Task:",
        BODY_X, sep_y + Inches(0.10),
        Inches(2.0), Inches(0.35),
        size=13, bold=True, color=DEEP_INDIGO)

    data_items = [
        "Response time (ms)",
        "Selection accuracy",
        "Attempt count",
        "Pointer trajectory",
        "Break frequency",
    ]
    col_w = BODY_W / len(data_items)
    cx = BODY_X
    for item in data_items:
        add_rect(slide, cx + Inches(0.04), sep_y + Inches(0.50),
                 col_w - Inches(0.08), Inches(0.55), INDIGO_50, INDIGO_200, Pt(0.5))
        txb(slide, item, cx + Inches(0.08), sep_y + Inches(0.55),
            col_w - Inches(0.16), Inches(0.45),
            size=11.5, color=INDIGO_700, bold=True, align=PP_ALIGN.CENTER)
        cx += col_w

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — GAME DESIGN
# ═══════════════════════════════════════════════════════════════════════════
def slide_game_design(prs):
    slide = make_base_slide(prs, "Game Design — Six Gameplay Chapters", 6)

    chapters = [
        ("Ch. 1", "My World",
         "Baseline",
         "Avatar setup · Name response · Joint attention",
         INDIGO_50, INDIGO_700),
        ("Ch. 2", "Feeling World",
         "Social Communication (40%)",
         "Emotion matching (drag-drop) · Expression mirror · Regulation scenarios",
         rgb(0xFF,0xF7,0xED), rgb(0xD9,0x77,0x06)),
        ("Ch. 3", "Social World",
         "Social Communication (40%)",
         "Greeting · Conversation + ToM probe · Sharing · Imitation (Copy Cat)",
         rgb(0xEF,0xF6,0xFF), rgb(0x1D,0x6F,0xBF)),
        ("Ch. 4", "Routine & Patterns",
         "Restricted/Repetitive (30%)",
         "Routine sequencing · Flexibility challenge · Pattern detective · Special interests",
         rgb(0xF0,0xFD,0xF4), rgb(0x04,0x78,0x57)),
        ("Ch. 5", "Pretend & Senses",
         "Pretend Play (15%) + Sensory (15%)",
         "Pretend recognition · Create pretend world · Sensory explorer (Tone.js sounds)",
         rgb(0xFA,0xF5,0xFF), rgb(0x7C,0x3A,0xED)),
        ("Ch. 6", "Grand Finale",
         "Consistency Check",
         "15 sampled tasks from Ch. 1–5 · Response consistency verification",
         rgb(0xFF,0xFB,0xEB), rgb(0xD9,0x77,0x06)),
    ]

    row_h = Inches(0.57)
    row_y = BODY_Y_START + Inches(0.05)
    col_widths = [Inches(0.65), Inches(1.35), Inches(1.85), Inches(3.78)]
    headers = ["", "Chapter", "Domain (Weight)", "Key Tasks"]

    # Header row
    cx = BODY_X
    for j, hdr in enumerate(headers):
        add_rect(slide, cx, row_y, col_widths[j], Inches(0.38), DEEP_INDIGO)
        txb(slide, hdr, cx + Inches(0.05), row_y + Inches(0.05),
            col_widths[j] - Inches(0.1), Inches(0.28),
            size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += col_widths[j]

    # Chapter rows
    for i, (ch_num, ch_name, domain, tasks, bg, col) in enumerate(chapters):
        ry = row_y + Inches(0.38) + row_h * i
        cx = BODY_X

        # Colored tag cell
        add_rect(slide, cx, ry, col_widths[0], row_h, col)
        txb(slide, ch_num, cx + Inches(0.02), ry + Inches(0.1),
            col_widths[0] - Inches(0.04), row_h - Inches(0.1),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += col_widths[0]

        # Chapter name
        add_rect(slide, cx, ry, col_widths[1], row_h, bg)
        txb(slide, ch_name, cx + Inches(0.06), ry + Inches(0.08),
            col_widths[1] - Inches(0.12), row_h - Inches(0.08),
            size=12.5, bold=True, color=col, align=PP_ALIGN.LEFT)
        cx += col_widths[1]

        # Domain
        add_rect(slide, cx, ry, col_widths[2], row_h,
                 LIGHT_GRAY if i % 2 == 0 else OFF_WHITE)
        txb(slide, domain, cx + Inches(0.06), ry + Inches(0.1),
            col_widths[2] - Inches(0.12), row_h - Inches(0.1),
            size=10.5, color=SLATE_800, align=PP_ALIGN.LEFT)
        cx += col_widths[2]

        # Tasks
        add_rect(slide, cx, ry, col_widths[3], row_h,
                 LIGHT_GRAY if i % 2 == 0 else OFF_WHITE)
        txb(slide, tasks, cx + Inches(0.06), ry + Inches(0.08),
            col_widths[3] - Inches(0.12), row_h - Inches(0.08),
            size=10.5, color=SLATE_600, align=PP_ALIGN.LEFT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — SCORING ENGINE
# ═══════════════════════════════════════════════════════════════════════════
def slide_scoring(prs):
    slide = make_base_slide(prs, "Scoring Engine — Research-Grounded Design", 7)

    # Three source papers - left column
    col_left_w = Inches(3.55)
    col_right_w = BODY_W - col_left_w - Inches(0.25)
    col2_x = BODY_X + col_left_w + Inches(0.25)

    # Left: Source frameworks
    add_rect(slide, BODY_X, BODY_Y_START, col_left_w, Inches(0.38), INDIGO_100)
    txb(slide, "Three Research Foundations",
        BODY_X + Inches(0.1), BODY_Y_START + Inches(0.05),
        col_left_w - Inches(0.2), Inches(0.28),
        size=13.5, bold=True, color=DEEP_INDIGO)

    sources = [
        ("EmoGalaxy  (Irani et al., 2018)",
         "Task accuracy = correct_answers / total_moves",
         INDIGO_50, INDIGO_700),
        ("DTT Game  (Khowaja & Salim, 2018)",
         "Avg. attempts = Σ attempts / total_questions",
         GREEN_100, GREEN_700),
        ("ADOS-2  (Maddox et al., 2017)",
         "Domain weighting: SC=0.40, RRB=0.30, PP=0.15, SP=0.15",
         AMBER_100, AMBER_600),
    ]

    sy = BODY_Y_START + Inches(0.42)
    for title, formula, bg, col in sources:
        add_rect(slide, BODY_X, sy, col_left_w, Inches(0.82), bg, col, Pt(0.75))
        txb(slide, title,
            BODY_X + Inches(0.12), sy + Inches(0.07),
            col_left_w - Inches(0.24), Inches(0.33),
            size=12.5, bold=True, color=col)
        txb(slide, formula,
            BODY_X + Inches(0.12), sy + Inches(0.40),
            col_left_w - Inches(0.24), Inches(0.36),
            size=11, italic=True, color=SLATE_600)
        sy += Inches(0.90)

    # Right column: Combined formula + domain table
    add_rect(slide, col2_x, BODY_Y_START, col_right_w, Inches(0.38), INDIGO_100)
    txb(slide, "Combined Weighted Score",
        col2_x + Inches(0.1), BODY_Y_START + Inches(0.05),
        col_right_w - Inches(0.2), Inches(0.28),
        size=13.5, bold=True, color=DEEP_INDIGO)

    # Formula box
    formula_y = BODY_Y_START + Inches(0.46)
    add_rect(slide, col2_x, formula_y, col_right_w, Inches(0.72),
             DEEP_INDIGO)
    txb(slide, "S_combined = Σ  w_d · R_d  ×  min( Π m_f ,  2.0 )",
        col2_x + Inches(0.12), formula_y + Inches(0.12),
        col_right_w - Inches(0.24), Inches(0.52),
        size=13.5, bold=True, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    txb(slide, "w_d = domain weight   ·   R_d = raw domain score   ·   m_f = red-flag multiplier (capped at 2.0×)",
        col2_x, formula_y + Inches(0.78),
        col_right_w, Inches(0.30),
        size=9.5, italic=True, color=SLATE_600, align=PP_ALIGN.CENTER)

    # Domain table
    dt_y = formula_y + Inches(1.15)
    dt_headers = ["Domain", "Chapters", "Weight", "High Risk ≥"]
    dt_rows = [
        ["Social Communication", "2, 3", "0.40", "46 pts"],
        ["Restricted / Repetitive", "4", "0.30", "31 pts"],
        ["Pretend Play", "5 (L1–L2)", "0.15", "21 pts"],
        ["Sensory Processing", "5 (L3)", "0.15", "16 pts"],
    ]
    col_ws = [Inches(1.95), Inches(0.85), Inches(0.68), Inches(0.90)]
    rh = Inches(0.35)
    cx = col2_x
    for j, hdr in enumerate(dt_headers):
        add_rect(slide, cx, dt_y, col_ws[j], rh, DEEP_INDIGO)
        txb(slide, hdr, cx + Inches(0.04), dt_y + Inches(0.05),
            col_ws[j] - Inches(0.08), rh - Inches(0.05),
            size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += col_ws[j]

    for i, row in enumerate(dt_rows):
        cy = dt_y + rh * (i + 1)
        cx = col2_x
        for j, cell in enumerate(row):
            bg = INDIGO_50 if i % 2 == 0 else OFF_WHITE
            add_rect(slide, cx, cy, col_ws[j], rh, bg)
            txb(slide, cell, cx + Inches(0.04), cy + Inches(0.06),
                col_ws[j] - Inches(0.08), rh - Inches(0.06),
                size=10.5, color=SLATE_800, align=PP_ALIGN.CENTER)
            cx += col_ws[j]

    # Vertical divider
    add_rect(slide, BODY_X + col_left_w + Inches(0.10), BODY_Y_START,
             Inches(0.025), BODY_H, INDIGO_200)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — RED FLAG DETECTION
# ═══════════════════════════════════════════════════════════════════════════
def slide_red_flags(prs):
    slide = make_base_slide(prs, "Automated Red-Flag Detection", 8)

    # Intro text
    txb(slide, "Five clinically-grounded behavioural red flags trigger multiplicative severity adjustments.",
        BODY_X, BODY_Y_START, BODY_W, Inches(0.38),
        size=14, italic=True, color=SLATE_600)

    flags = [
        ("×1.20", "Negative Emotion Recognition < 50%",
         "Ch. 2 Level 1: accuracy on Sad + Scared cards below 50%\n→ Heightened social-communication concern"),
        ("×1.30", "Complete Absence of Pretend Play",
         "Ch. 5 Level 1: all 5 pretend scenarios answered literally\n→ Strongest single red-flag multiplier"),
        ("×1.15", "Extreme Sensory Distress",
         "Ch. 5 Level 3: ≥ 4 of 6 Tone.js-synthesised sounds rated distressing\n→ Atypical sensory reactivity"),
        ("×1.20", "Rigid Pattern Distress",
         "Ch. 4 Level 3: meltdown at forced error + insists on returning to original pattern\n→ RRB indicator"),
        ("×1.25", "Poor Imitation Across All Modalities",
         "Ch. 3 Level 4: ≥ 5 errors across facial, body, and object imitation tasks\n→ Multi-modal imitation deficit"),
    ]

    card_h = Inches(0.68)
    gap = Inches(0.10)
    fy = BODY_Y_START + Inches(0.44)
    mult_w = Inches(0.72)
    name_w = Inches(2.20)
    desc_w = BODY_W - mult_w - name_w - Inches(0.1)

    for mult, name, desc in flags:
        # Multiplier badge
        add_rect(slide, BODY_X, fy, mult_w, card_h, RED_600)
        txb(slide, mult, BODY_X + Inches(0.04), fy + Inches(0.12),
            mult_w - Inches(0.08), card_h - Inches(0.12),
            size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Flag name
        add_rect(slide, BODY_X + mult_w, fy, name_w, card_h, AMBER_100, AMBER_600, Pt(0.5))
        txb(slide, name,
            BODY_X + mult_w + Inches(0.1), fy + Inches(0.1),
            name_w - Inches(0.2), card_h - Inches(0.1),
            size=12, bold=True, color=rgb(0x78,0x35,0x00), wrap=True)

        # Description
        add_rect(slide, BODY_X + mult_w + name_w, fy,
                 desc_w + Inches(0.1), card_h, LIGHT_GRAY)
        txb(slide, desc,
            BODY_X + mult_w + name_w + Inches(0.1), fy + Inches(0.06),
            desc_w - Inches(0.1), card_h - Inches(0.06),
            size=10.5, color=SLATE_800, wrap=True)

        fy += card_h + gap

    # Cap note
    note_y = fy + Inches(0.05)
    add_rect(slide, BODY_X, note_y, BODY_W, Inches(0.40), INDIGO_50, INDIGO_200, Pt(0.5))
    txb(slide, "Multipliers stack multiplicatively and are capped at 2.0× to prevent extreme score inflation.  "
               "Each flag is independently verified by integration tests using simulated full-session data.",
        BODY_X + Inches(0.15), note_y + Inches(0.07),
        BODY_W - Inches(0.3), Inches(0.30),
        size=11, italic=True, color=SLATE_600)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — IMPLEMENTATION RESULTS
# ═══════════════════════════════════════════════════════════════════════════
def slide_results(prs):
    slide = make_base_slide(prs, "Implementation Results", 9)

    # Stat cards row
    stats = [
        ("312", "Automated tests\n(all passing)"),
        ("5 / 5", "Objectives\nachieved"),
        ("6", "Game chapters\nimplemented"),
        ("5", "Red flags\nverified"),
    ]
    card_w = (BODY_W - Inches(0.45)) / 4
    card_h = Inches(0.95)
    cx = BODY_X
    for val, lbl in stats:
        add_rect(slide, cx, BODY_Y_START, card_w, card_h, INDIGO_50, INDIGO_200, Pt(0.75))
        txb(slide, val, cx + Inches(0.06), BODY_Y_START + Inches(0.06),
            card_w - Inches(0.12), Inches(0.50),
            size=28, bold=True, color=INDIGO_700, align=PP_ALIGN.CENTER)
        txb(slide, lbl, cx + Inches(0.06),
            BODY_Y_START + Inches(0.56),
            card_w - Inches(0.12), Inches(0.38),
            size=10, color=SLATE_600, align=PP_ALIGN.CENTER)
        cx += card_w + Inches(0.15)

    # Objectives achieved
    obj_y = BODY_Y_START + card_h + Inches(0.18)
    txb(slide, "All Objectives Delivered:",
        BODY_X, obj_y, BODY_W, Inches(0.32),
        size=14, bold=True, color=DEEP_INDIGO)

    obj_y += Inches(0.34)
    objectives = [
        "Six gameplay chapters covering all four ADOS-2 / DSM-5 domains, deployed to Vercel",
        "Penalty-based scoring engine (EmoGalaxy + DTT + ADOS-2 formulae) with complete unit test coverage",
        "Five red flags with multiplicative weighting — verified in integration tests with simulated session data",
        "Researcher dashboard: Better Auth login · domain radar charts · pointer heatmaps · JSON/CSV export",
        "HMAC-SHA256 signed caregiver report tokens generated at session completion — correctness verified",
    ]

    col_w = (BODY_W - Inches(0.15)) / 2
    col2_x = BODY_X + col_w + Inches(0.15)
    left_items = objectives[:3]
    right_items = objectives[3:]

    cur_y = obj_y
    for item in left_items:
        txb(slide, "✓", BODY_X, cur_y, Inches(0.28), Inches(0.44),
            size=14, bold=True, color=GREEN_700)
        txb(slide, item, BODY_X + Inches(0.28), cur_y,
            col_w - Inches(0.28), Inches(0.46),
            size=12, color=SLATE_800, wrap=True)
        cur_y += Inches(0.50)

    cur_y = obj_y
    for item in right_items:
        txb(slide, "✓", col2_x, cur_y, Inches(0.28), Inches(0.44),
            size=14, bold=True, color=GREEN_700)
        txb(slide, item, col2_x + Inches(0.28), cur_y,
            col_w - Inches(0.28), Inches(0.46),
            size=12, color=SLATE_800, wrap=True)
        cur_y += Inches(0.50)

    # Tech note
    note_y = BODY_Y_END - Inches(0.38)
    txb(slide, "Serverless PostgreSQL (Neon) · Zero audio files shipped (Tone.js synthesis at runtime) · "
               "Zero image files (emoji + CSS + SVG only)  ·  Production-deployed on Vercel free tier",
        BODY_X, note_y, BODY_W, Inches(0.35),
        size=10, italic=True, color=SLATE_600)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — ETHICAL GUARDRAILS & IMPACT
# ═══════════════════════════════════════════════════════════════════════════
def slide_ethics(prs):
    slide = make_base_slide(prs, "Ethical Guardrails & Societal Impact", 10)

    col_left_w  = Inches(3.85)
    col_right_w = BODY_W - col_left_w - Inches(0.25)
    col2_x = BODY_X + col_left_w + Inches(0.25)

    # Left: ethical guardrails
    add_rect(slide, BODY_X, BODY_Y_START, col_left_w, Inches(0.38), INDIGO_100)
    txb(slide, "Built-in Ethical Guardrails",
        BODY_X + Inches(0.1), BODY_Y_START + Inches(0.05),
        col_left_w - Inches(0.2), Inches(0.28),
        size=13.5, bold=True, color=DEEP_INDIGO)

    ethics_items = [
        ("Non-diagnostic disclaimer",
         "Shown on every results screen and caregiver report — this tool screens, not diagnoses"),
        ("Data minimisation (GDPR-aligned)",
         "Session ID is the only permanent identifier; player name is an optional display label only"),
        ("No biometrics collected",
         "No photos, video, voice recordings, or precise location data at any point"),
        ("HMAC-SHA256 signed reports",
         "Caregiver links are tamper-proof; unsigned tokens return an error page"),
        ("Safety: no flicker, no harsh audio",
         "Framer Motion < 3 Hz (IEC 61966); Tone.js attenuated −12 dB / −6 dB in low/medium mode"),
    ]

    ey = BODY_Y_START + Inches(0.44)
    for title, desc in ethics_items:
        txb(slide, "▸  " + title,
            BODY_X, ey, col_left_w, Inches(0.30),
            size=12, bold=True, color=INDIGO_700)
        txb(slide, desc,
            BODY_X + Inches(0.15), ey + Inches(0.30),
            col_left_w - Inches(0.15), Inches(0.28),
            size=10.5, italic=True, color=SLATE_600, wrap=True)
        ey += Inches(0.62)

    # Right: societal impact
    add_rect(slide, col2_x, BODY_Y_START, col_right_w, Inches(0.38), INDIGO_100)
    txb(slide, "Societal Impact",
        col2_x + Inches(0.1), BODY_Y_START + Inches(0.05),
        col_right_w - Inches(0.2), Inches(0.28),
        size=13.5, bold=True, color=DEEP_INDIGO)

    impact_items = [
        ("Clinical pre-referral screening",
         "Deploy on any browser device — HMAC report forwarded directly to the specialist"),
        ("Telehealth-compatible",
         "Zero-install, zero-hardware — suitable where face-to-face assessment is unavailable"),
        ("Parental awareness",
         "Domain-level risk summaries guide more informed referral conversations with physicians"),
        ("Population research",
         "Anonymous bulk export enables epidemiological studies where ADOS-2 is logistically impractical"),
    ]

    iy = BODY_Y_START + Inches(0.44)
    for title, desc in impact_items:
        txb(slide, "▸  " + title,
            col2_x, iy, col_right_w, Inches(0.30),
            size=12, bold=True, color=GREEN_700)
        txb(slide, desc,
            col2_x + Inches(0.15), iy + Inches(0.30),
            col_right_w - Inches(0.15), Inches(0.28),
            size=10.5, italic=True, color=SLATE_600, wrap=True)
        iy += Inches(0.62)

    # Vertical divider
    add_rect(slide, BODY_X + col_left_w + Inches(0.10), BODY_Y_START,
             Inches(0.025), BODY_H, INDIGO_200)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — CONCLUSION & FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    slide = make_base_slide(prs, "Conclusion & Future Work", 11)

    col_left_w  = Inches(3.85)
    col_right_w = BODY_W - col_left_w - Inches(0.25)
    col2_x = BODY_X + col_left_w + Inches(0.25)

    # Left: conclusion
    add_rect(slide, BODY_X, BODY_Y_START, col_left_w, Inches(0.38), INDIGO_100)
    txb(slide, "What We Achieved",
        BODY_X + Inches(0.1), BODY_Y_START + Inches(0.05),
        col_left_w - Inches(0.2), Inches(0.28),
        size=13.5, bold=True, color=DEEP_INDIGO)

    conclusions = [
        "First open-source, browser-based, multi-domain ASD serious game covering all four ADOS-2 / DSM-5 domains in a single 60–75 min session",
        "Research-grounded penalty scoring engine combining three independent peer-reviewed frameworks",
        "Automated red-flag detection with multiplicative severity weighting — verified by 312 automated tests",
        "Production-ready deployment: serverless PostgreSQL, zero static assets, free-tier Vercel hosting",
        "HMAC-signed caregiver report shareable without a researcher account",
    ]

    cy = BODY_Y_START + Inches(0.44)
    for item in conclusions:
        txb(slide, "▸", BODY_X, cy, Inches(0.26), Inches(0.50),
            size=14, bold=True, color=INDIGO_700)
        txb(slide, item, BODY_X + Inches(0.26), cy,
            col_left_w - Inches(0.26), Inches(0.54),
            size=12, color=SLATE_800, wrap=True)
        cy += Inches(0.58)

    # Right: future work
    add_rect(slide, col2_x, BODY_Y_START, col_right_w, Inches(0.38), INDIGO_100)
    txb(slide, "Future Directions",
        col2_x + Inches(0.1), BODY_Y_START + Inches(0.05),
        col_right_w - Inches(0.2), Inches(0.28),
        size=13.5, bold=True, color=DEEP_INDIGO)

    future = [
        "Large-scale clinical validation (200+ participants, diverse age / culture subgroups)",
        "ML-enhanced scoring — SVM / gradient-boosted classifier on game-response feature vectors, targeting ≥ 93% EmoGalaxy accuracy",
        "Age-stratified norms — separate scoring coefficients for ages 3–5, 6–8, and 9–10",
        "Multilingual & cultural adaptation — stimulus text and social scenarios for non-English populations",
        "Integration with referral workflows — HMAC report piped into clinic patient management systems",
    ]

    fy = BODY_Y_START + Inches(0.44)
    for item in future:
        txb(slide, "▸", col2_x, fy, Inches(0.26), Inches(0.50),
            size=14, bold=True, color=AMBER_600)
        txb(slide, item, col2_x + Inches(0.26), fy,
            col_right_w - Inches(0.26), Inches(0.54),
            size=12, color=SLATE_800, wrap=True)
        fy += Inches(0.58)

    # Vertical divider
    add_rect(slide, BODY_X + col_left_w + Inches(0.10), BODY_Y_START,
             Inches(0.025), BODY_H, INDIGO_200)

    # Disclaimer banner
    disc_y = BODY_Y_END - Inches(0.15)
    # Already close to footer — just ensure it's within body
    # (banner rendered inside slide body before the chrome footer)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — THANK YOU / Q&A
# ═══════════════════════════════════════════════════════════════════════════
def slide_thankyou(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    add_rect(slide, 0, 0, W, H, OFF_WHITE)
    add_rect(slide, 0, 0, Inches(0.45), H, INDIGO_700)
    add_rect(slide, 0, H - Inches(0.70), W, Inches(0.70), DEEP_INDIGO)
    add_rect(slide, 0, H - Inches(0.70), Inches(0.45), Inches(0.70), INDIGO_700)

    # Divider
    add_rect(slide, Inches(0.65), Inches(2.52), W - Inches(0.9), Inches(0.025), INDIGO_200)

    txb(slide, "Thank you",
        Inches(0.65), Inches(0.80),
        W - Inches(0.9), Inches(1.10),
        size=52, bold=True, color=INDIGO_700, align=PP_ALIGN.LEFT)

    txb(slide, "Questions & Discussion",
        Inches(0.65), Inches(1.92),
        W - Inches(0.9), Inches(0.55),
        size=20, color=SLATE_800, align=PP_ALIGN.LEFT)

    txb(slide,
        "Md. Shifat Hasan (2107067)   ·   Abu Hasanat Soykot (2107100)\n"
        "Supervisor: Dr. Md. Aminul Haque Akhand, Professor, Dept. of CSE",
        Inches(0.65), Inches(2.65),
        W - Inches(0.9), Inches(0.70),
        size=13.5, bold=True, color=DEEP_INDIGO, align=PP_ALIGN.LEFT)

    txb(slide,
        "Dept. of Computer Science and Engineering\n"
        "Khulna University of Engineering & Technology  ·  Khulna 9203, Bangladesh\n"
        "April 2026",
        Inches(0.65), Inches(3.40),
        W - Inches(0.9), Inches(0.75),
        size=12, color=SLATE_600, align=PP_ALIGN.LEFT)

    # Non-diagnostic reminder
    nd_y = H - Inches(1.40)
    add_rect(slide, Inches(0.65), nd_y, W - Inches(0.9), Inches(0.50),
             AMBER_100, AMBER_600, Pt(0.75))
    txb(slide,
        "Horizons is a screening guide only — it does not make a diagnosis. "
        "Always consult a qualified healthcare specialist.",
        Inches(0.80), nd_y + Inches(0.09),
        W - Inches(1.2), Inches(0.38),
        size=11, italic=True, color=rgb(0x78, 0x35, 0x00), align=PP_ALIGN.CENTER)

    txb(slide, FOOTER_TEXT,
        Inches(0.65), H - Inches(0.57),
        W - Inches(0.9), Inches(0.42),
        size=9, color=LAVENDER, align=PP_ALIGN.LEFT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)        # 1 (no number shown)
    slide_outline(prs)      # 1/12
    slide_problem(prs)      # 2/12
    slide_related_work(prs) # 3/12
    slide_overview(prs)     # 4/12
    slide_architecture(prs) # 5/12
    slide_game_design(prs)  # 6/12
    slide_scoring(prs)      # 7/12
    slide_red_flags(prs)    # 8/12
    slide_results(prs)      # 9/12
    slide_ethics(prs)       # 10/12
    slide_conclusion(prs)   # 11/12
    slide_thankyou(prs)     # 12 (no number shown)

    out = "/Users/shifathasan/Developer/Academic/_PROJECTS_/SYSTEM/horizons/docs/slides/slides.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
